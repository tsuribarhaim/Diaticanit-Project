from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


def add_title_bar(slide, text):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_slide(prs, section_title, headline, bullets, note):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    add_title_bar(slide, "Personal Health Companion")

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.9), Inches(0.9))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = section_title
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(30, 41, 59)

    headline_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(2.0),
        Inches(11.9),
        Inches(0.95),
    )
    headline_box.fill.solid()
    headline_box.fill.fore_color.rgb = RGBColor(220, 252, 231)
    headline_box.line.color.rgb = RGBColor(74, 222, 128)

    headline_tf = headline_box.text_frame
    headline_tf.clear()
    headline_p = headline_tf.paragraphs[0]
    headline_run = headline_p.add_run()
    headline_run.text = headline
    headline_run.font.name = "Calibri"
    headline_run.font.size = Pt(18)
    headline_run.font.bold = True
    headline_run.font.color.rgb = RGBColor(22, 101, 52)

    bullets_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.15), Inches(11.4), Inches(3.8))
    bullets_tf = bullets_box.text_frame
    bullets_tf.word_wrap = True
    bullets_tf.clear()

    for i, bullet in enumerate(bullets):
        p = bullets_tf.paragraphs[0] if i == 0 else bullets_tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 41, 59)

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = f"Speaker cue: {note}"


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        {
            "section_title": "Mission and Outcome",
            "headline": "From Data Capture to Actionable Health Decisions",
            "bullets": [
                "Mission: convert fragmented personal health inputs into secure, structured, and usable insights.",
                "Current outcome: users can onboard, upload records, set goals, and track daily behavior in one workflow.",
                "Delivery status: core platform is live across Phases 1-4 with controlled AI augmentation.",
            ],
            "note": "Lead with value delivered, not architecture.",
        },
        {
            "section_title": "Why This Matters",
            "headline": "Business Value",
            "bullets": [
                "Before: health data was scattered, hard to interpret, and easy to ignore.",
                "Now: the app creates one user-owned flow for profile, documents, goals, and daily progress.",
                "Result: higher adherence potential, clearer progress visibility, and a stronger personalization foundation.",
            ],
            "note": "Frame as reduction in friction and increase in user trust.",
        },
        {
            "section_title": "What Is Live Today",
            "headline": "Production-Ready Capabilities",
            "bullets": [
                "Secure auth, protected routes, and owner-scoped data access.",
                "Profile onboarding and document extraction review workflows.",
                "Goals model and daily report tracking with reusable defaults.",
                "Confidence/confirmation gating and AI retry path for parse quality.",
            ],
            "note": "This is the headline slide for stakeholders asking what users can do now.",
        },
        {
            "section_title": "Delivery by Phase",
            "headline": "Execution Progress",
            "bullets": [
                "Phase 1: secure foundation and critical user journeys.",
                "Phase 2: extraction pipeline with deterministic baseline and AI-assisted path.",
                "Phase 3: goals and progress scaffolding for measurable outcomes.",
                "Phase 4: daily reporting maturity, UX polish, and parse-mode transparency.",
            ],
            "note": "Keep this concise and momentum-focused.",
        },
        {
            "section_title": "System Flow",
            "headline": "How Value Is Produced End-to-End",
            "bullets": [
                "1. User authenticates and completes profile baseline.",
                "2. User provides inputs through documents and daily reports.",
                "3. System parses inputs via heuristic default or optional AI mode.",
                "4. Confidence and confirmation gates protect data quality before progress impact.",
                "5. Structured outputs connect into goals and progress views.",
            ],
            "note": "Emphasize human-in-the-loop trust design.",
        },
        {
            "section_title": "Risk and Control Posture",
            "headline": "Top Risks and Mitigations",
            "bullets": [
                "AI config risk is mitigated by deterministic fallback and explicit mode controls.",
                "Parse trust risk is mitigated by confidence scoring and confirmation gates.",
                "Environment drift risk is mitigated by migration discipline and deployment runbook.",
            ],
            "note": "Show that risk is managed by design.",
        },
        {
            "section_title": "Next 30/60 Days",
            "headline": "Execution Plan",
            "bullets": [
                "30 days: harden preview/prod rollout and add parse quality telemetry.",
                "30 days: expand regression tests for mode switching and daily-report reliability.",
                "60 days: improve AI operational defaults and add outcome KPI dashboards.",
                "60 days: run structured UAT and prepare release signoff package.",
            ],
            "note": "Focus on quality scaling and adoption measurement.",
        },
        {
            "section_title": "Decisions and Asks",
            "headline": "Leadership Alignment Needed",
            "bullets": [
                "Approve reliability and telemetry as next release priority.",
                "Confirm preview-first deployment gate before production promotion.",
                "Align on success KPIs: active reporting rate, confirmation rate, AI retry success, and time-to-insight.",
            ],
            "note": "End with specific decisions to maintain momentum.",
        },
    ]

    for slide_data in slides:
        add_slide(
            prs=prs,
            section_title=slide_data["section_title"],
            headline=slide_data["headline"],
            bullets=slide_data["bullets"],
            note=slide_data["note"],
        )

    prs.save(output_path)


if __name__ == "__main__":
    output = Path(__file__).resolve().parent / "executive-powerpoint-report-board-8-slides.pptx"
    build_presentation(output)
    print(f"Created: {output}")
